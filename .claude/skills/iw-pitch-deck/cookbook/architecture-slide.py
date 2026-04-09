# /// script
# name = "architecture-slide"
# purpose = "Full-width diagram slide for architecture overviews, system diagrams, or data flows"
# best_for = "System architecture, deployment topology, data flow, C4 diagrams"
# avoid_when = "Simple bullet points (use content-slide), metrics (use stats-slide)"
# instructions = "Replace BRAND_* variables with brand.json values. Render Mermaid diagram to PNG first: mmdc -i file.mmd -o file.png -b white -w 1920. Then use slide.shapes.add_picture() to embed."
# diagram_format = "PNG (not SVG — PPTX needs raster images)"
# requires = ["python-pptx", "Pillow"]
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BRAND_PRIMARY = "#1B2A4A"
BRAND_SECONDARY = "#2E86AB"
BRAND_ACCENT = "#F18F01"
BRAND_BG = "#FFFFFF"
BRAND_BG_ALT = "#F5F7FA"
BRAND_TEXT = "#1A1A2E"
BRAND_TEXT_LIGHT = "#6B7280"
BRAND_FONT_HEADING = "Inter"
BRAND_FONT_BODY = "Inter"


def hex_to_rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Background ───────────────────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

# ── Slide title ──────────────────────────────────────────────────────
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.333), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "System Architecture"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = hex_to_rgb(BRAND_PRIMARY)
p.font.name = BRAND_FONT_HEADING
p.alignment = PP_ALIGN.LEFT

# ── Title underline ──────────────────────────────────────────────────
line = slide.shapes.add_shape(1, Inches(1.0), Inches(1.2), Inches(2.0), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
line.line.fill.background()

# ── Diagram area ─────────────────────────────────────────────────────
# In production, replace this placeholder with:
#   slide.shapes.add_picture("diagram.png", left, top, width, height)
#
# Rendering workflow:
#   1. mmdc -i diagram.mmd -o diagram.png -b white -w 1920
#   2. slide.shapes.add_picture("diagram.png", Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.4))

diagram_area = slide.shapes.add_shape(1, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.4))
diagram_area.fill.solid()
diagram_area.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
diagram_area.line.color.rgb = hex_to_rgb("#E5E7EB")
diagram_area.line.width = Pt(1)

# Placeholder label
label = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.733), Inches(1.0))
label_tf = label.text_frame
label_tf.word_wrap = True
label_p = label_tf.paragraphs[0]
label_p.text = "[ Architecture Diagram — render .mmd to .png and embed here ]"
label_p.font.size = Pt(16)
label_p.font.color.rgb = hex_to_rgb(BRAND_TEXT_LIGHT)
label_p.font.name = BRAND_FONT_BODY
label_p.alignment = PP_ALIGN.CENTER

# ── Caption ──────────────────────────────────────────────────────────
caption = slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(11.733), Inches(0.3))
cap_tf = caption.text_frame
cap_p = cap_tf.paragraphs[0]
cap_p.text = "Figure 1: High-level system architecture"
cap_p.font.size = Pt(11)
cap_p.font.italic = True
cap_p.font.color.rgb = hex_to_rgb(BRAND_TEXT_LIGHT)
cap_p.font.name = BRAND_FONT_BODY
cap_p.alignment = PP_ALIGN.CENTER

# ── Speaker notes ────────────────────────────────────────────────────
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = "Speaker notes: Walk through each component in the diagram from left to right (or top to bottom)."

prs.save("architecture-slide.pptx")
print("Created: architecture-slide.pptx")
