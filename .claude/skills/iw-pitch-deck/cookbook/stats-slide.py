# /// script
# name = "stats-slide"
# purpose = "KPI/metrics slide with 2-4 large numbers and labels in a horizontal layout"
# best_for = "Cost savings, performance metrics, growth numbers, before/after comparisons"
# avoid_when = "Qualitative content (use content-slide), narrative points (use quote-slide)"
# instructions = "Replace BRAND_* variables with brand.json values. Use 2-4 stats max. Each stat has a large number and a short label (under 5 words). Numbers should be bold and eye-catching."
# max_stats = 4
# requires = ["python-pptx"]
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BRAND_PRIMARY = "#1B2A4A"
BRAND_SECONDARY = "#2E86AB"
BRAND_ACCENT = "#F18F01"
BRAND_BG = "#FFFFFF"
BRAND_TEXT = "#1A1A2E"
BRAND_TEXT_LIGHT = "#6B7280"
BRAND_FONT_HEADING = "Inter"
BRAND_FONT_BODY = "Inter"


def hex_to_rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Background ───────────────────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

# ── Slide title ──────────────────────────────────────────────────────
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Key Metrics"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = hex_to_rgb(BRAND_PRIMARY)
p.font.name = BRAND_FONT_HEADING
p.alignment = PP_ALIGN.LEFT

# ── Title underline ──────────────────────────────────────────────────
line = slide.shapes.add_shape(1, Inches(1.0), Inches(1.55), Inches(2.0), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
line.line.fill.background()

# ── Stats data ───────────────────────────────────────────────────────
stats = [
    ("97%", "Cost Reduction"),
    ("30x", "Faster Processing"),
    ("24/7", "Automated Operation"),
]

stat_count = len(stats)
total_width = 11.333
stat_width = total_width / stat_count
start_x = 1.0
start_y = 2.5

for i, (number, label) in enumerate(stats):
    x = start_x + i * stat_width

    # Accent dot above number
    dot = slide.shapes.add_shape(
        1, Inches(x + stat_width / 2 - 0.15), Inches(start_y), Inches(0.3), Inches(0.06)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    dot.line.fill.background()

    # Large number
    num_box = slide.shapes.add_textbox(
        Inches(x), Inches(start_y + 0.3), Inches(stat_width), Inches(1.5)
    )
    num_tf = num_box.text_frame
    num_tf.word_wrap = True
    num_p = num_tf.paragraphs[0]
    num_p.text = number
    num_p.font.size = Pt(54)
    num_p.font.bold = True
    num_p.font.color.rgb = hex_to_rgb(BRAND_SECONDARY)
    num_p.font.name = BRAND_FONT_HEADING
    num_p.alignment = PP_ALIGN.CENTER

    # Label
    lbl_box = slide.shapes.add_textbox(
        Inches(x), Inches(start_y + 2.0), Inches(stat_width), Inches(0.8)
    )
    lbl_tf = lbl_box.text_frame
    lbl_tf.word_wrap = True
    lbl_p = lbl_tf.paragraphs[0]
    lbl_p.text = label
    lbl_p.font.size = Pt(18)
    lbl_p.font.color.rgb = hex_to_rgb(BRAND_TEXT_LIGHT)
    lbl_p.font.name = BRAND_FONT_BODY
    lbl_p.alignment = PP_ALIGN.CENTER

# ── Divider lines between stats ──────────────────────────────────────
for i in range(1, stat_count):
    x = start_x + i * stat_width
    divider = slide.shapes.add_shape(
        1, Inches(x - 0.01), Inches(start_y + 0.2), Inches(0.02), Inches(2.8)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb("#E5E7EB")
    divider.line.fill.background()

# ── Speaker notes ────────────────────────────────────────────────────
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Speaker notes: Walk through each metric with context and comparison."
)

prs.save("stats-slide.pptx")
print("Created: stats-slide.pptx")
