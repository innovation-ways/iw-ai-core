# /// script
# name = "closing-slide"
# purpose = "Final slide with thank you message, CTA, and contact information"
# best_for = "Last slide of any presentation"
# avoid_when = "Never — every deck needs a closing slide"
# instructions = "Replace BRAND_* variables with brand.json values. Include company name, email, website, and social links from brand.json."
# requires = ["python-pptx"]
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BRAND_PRIMARY = "#1B2A4A"
BRAND_SECONDARY = "#2E86AB"
BRAND_ACCENT = "#F18F01"
BRAND_BG = "#FFFFFF"
BRAND_TEXT = "#1A1A2E"
BRAND_TEXT_LIGHT = "#6B7280"
BRAND_TEXT_ON_PRIMARY = "#FFFFFF"
BRAND_FONT_HEADING = "Inter"
BRAND_FONT_BODY = "Inter"
BRAND_COMPANY = "Innovation Ways"
BRAND_TAGLINE = "Engineering Tomorrow's Solutions"
BRAND_EMAIL = "info@innovationways.com"
BRAND_WEBSITE = "innovationways.com"
BRAND_LINKEDIN = "linkedin.com/company/innovation-ways"
BRAND_GITHUB = "github.com/innovation-ways"


def hex_to_rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Background ───────────────────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = hex_to_rgb(BRAND_PRIMARY)

# ── Left accent bar ─────────────────────────────────────────────────
bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
bar.fill.solid()
bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
bar.line.fill.background()

# ── Thank you / CTA ─────────────────────────────────────────────────
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = hex_to_rgb(BRAND_TEXT_ON_PRIMARY)
p.font.name = BRAND_FONT_HEADING
p.alignment = PP_ALIGN.LEFT

# ── CTA subtitle ─────────────────────────────────────────────────────
cta_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.333), Inches(0.8))
cta_tf = cta_box.text_frame
cta_tf.word_wrap = True
cta_p = cta_tf.paragraphs[0]
cta_p.text = "Let's discuss how we can help your project"
cta_p.font.size = Pt(22)
cta_p.font.color.rgb = hex_to_rgb(BRAND_SECONDARY)
cta_p.font.name = BRAND_FONT_BODY
cta_p.alignment = PP_ALIGN.LEFT

# ── Divider line ─────────────────────────────────────────────────────
divider = slide.shapes.add_shape(1, Inches(1.5), Inches(4.2), Inches(4.0), Inches(0.04))
divider.fill.solid()
divider.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
divider.line.fill.background()

# ── Contact details ──────────────────────────────────────────────────
contact_items = [
    f"\u2709  {BRAND_EMAIL}",
    f"\u2b24  {BRAND_WEBSITE}",
    f"\u2b24  {BRAND_LINKEDIN}",
    f"\u2b24  {BRAND_GITHUB}",
]

contact_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(10.333), Inches(2.5))
contact_tf = contact_box.text_frame
contact_tf.word_wrap = True

for i, item in enumerate(contact_items):
    if i == 0:
        p = contact_tf.paragraphs[0]
    else:
        p = contact_tf.add_paragraph()
    p.space_before = Pt(8)
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(200, 210, 225)
    p.font.name = BRAND_FONT_BODY
    p.alignment = PP_ALIGN.LEFT

# ── Company name at bottom ───────────────────────────────────────────
company_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.8), Inches(10.333), Inches(0.5))
co_tf = company_box.text_frame
co_p = co_tf.paragraphs[0]
co_p.text = f"{BRAND_COMPANY}  |  {BRAND_TAGLINE}"
co_p.font.size = Pt(12)
co_p.font.color.rgb = RGBColor(150, 160, 180)
co_p.font.name = BRAND_FONT_BODY
co_p.alignment = PP_ALIGN.LEFT

# ── Speaker notes ────────────────────────────────────────────────────
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = "Speaker notes: Thank the audience, open the floor for questions, and point to contact details on screen."

prs.save("closing-slide.pptx")
print("Created: closing-slide.pptx")
