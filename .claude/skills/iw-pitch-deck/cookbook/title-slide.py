# /// script
# name = "title-slide"
# purpose = "Opening slide with accent bar, centered title, subtitle, and date"
# best_for = "First slide of any presentation"
# avoid_when = "Never — every deck needs a title slide"
# instructions = "Replace BRAND_* variables with values from brand.json. Title should be 5-10 words. Subtitle is the tagline or one-line description."
# max_title_chars = 60
# max_subtitle_chars = 100
# requires = ["python-pptx"]
# ///

import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand placeholders (replace from brand.json) ─────────────────────
BRAND_PRIMARY = "#1B2A4A"
BRAND_SECONDARY = "#2E86AB"
BRAND_ACCENT = "#F18F01"
BRAND_BG = "#FFFFFF"
BRAND_TEXT = "#1A1A2E"
BRAND_TEXT_ON_PRIMARY = "#FFFFFF"
BRAND_FONT_HEADING = "Inter"
BRAND_FONT_BODY = "Inter"
BRAND_COMPANY = "Innovation Ways"
BRAND_TAGLINE = "Engineering Tomorrow's Solutions"


def hex_to_rgb(hex_color):
    """Convert hex color string to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Presentation setup ───────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# ── Background fill (prevents white-slide bug) ──────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = hex_to_rgb(BRAND_PRIMARY)

# ── Left accent bar ─────────────────────────────────────────────────
bar = slide.shapes.add_shape(
    1,  # MSO_SHAPE.RECTANGLE
    Inches(0),
    Inches(0),
    Inches(0.4),
    Inches(7.5),
)
bar.fill.solid()
bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
bar.line.fill.background()

# ── Title text ───────────────────────────────────────────────────────
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.333), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Presentation Title Here"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = hex_to_rgb(BRAND_TEXT_ON_PRIMARY)
p.font.name = BRAND_FONT_HEADING
p.alignment = PP_ALIGN.LEFT

# ── Subtitle text ────────────────────────────────────────────────────
txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.333), Inches(1.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "A concise subtitle describing the presentation"
p2.font.size = Pt(22)
p2.font.color.rgb = hex_to_rgb(BRAND_SECONDARY)
p2.font.name = BRAND_FONT_BODY
p2.alignment = PP_ALIGN.LEFT

# ── Date and company ─────────────────────────────────────────────────
txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(10.333), Inches(0.6))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = f"{BRAND_COMPANY}  |  {datetime.date.today().strftime('%B %Y')}"
p3.font.size = Pt(14)
p3.font.color.rgb = RGBColor(255, 255, 255)
p3.font.name = BRAND_FONT_BODY
p3.alignment = PP_ALIGN.LEFT

# ── Bottom accent line ───────────────────────────────────────────────
bottom_line = slide.shapes.add_shape(
    1,  # MSO_SHAPE.RECTANGLE
    Inches(1.5),
    Inches(6.8),
    Inches(4.0),
    Inches(0.05),
)
bottom_line.fill.solid()
bottom_line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
bottom_line.line.fill.background()

# ── Save ─────────────────────────────────────────────────────────────
prs.save("title-slide.pptx")
print("Created: title-slide.pptx")
