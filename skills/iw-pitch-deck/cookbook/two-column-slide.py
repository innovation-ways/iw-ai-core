# /// script
# name = "two-column-slide"
# purpose = "Split layout with text on the left and image/diagram area on the right"
# best_for = "Feature explanations with visuals, before/after comparisons, text + screenshot"
# avoid_when = "Full-width diagrams (use architecture-slide), pure text (use content-slide)"
# instructions = "Replace BRAND_* variables with brand.json values. Left column: title + 3-4 bullets or short paragraphs. Right column: image placeholder (use slide.shapes.add_picture() for real images)."
# max_left_words = 80
# requires = ["python-pptx"]
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BRAND_PRIMARY = "#1B2A4A"
BRAND_SECONDARY = "#2E86AB"
BRAND_ACCENT = "#F18F01"
BRAND_BG = "#FFFFFF"
BRAND_BG_ALT = "#F5F7FA"
BRAND_TEXT = "#1A1A2E"
BRAND_TEXT_LIGHT = "#6B7280"
BRAND_FONT_HEADING = "Inter"
BRAND_FONT_BODY = "Inter"


def hex_to_rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Background ───────────────────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

# ── Slide title (full width) ────────────────────────────────────────
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Feature Overview"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = hex_to_rgb(BRAND_PRIMARY)
p.font.name = BRAND_FONT_HEADING
p.alignment = PP_ALIGN.LEFT

# ── Title underline ──────────────────────────────────────────────────
line = slide.shapes.add_shape(1, Inches(1.0), Inches(1.55), Inches(2.0), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
line.line.fill.background()

# ── LEFT COLUMN: Text content ────────────────────────────────────────
left_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.5), Inches(4.5))
left_tf = left_box.text_frame
left_tf.word_wrap = True

bullets = [
    "Automated pipeline processes content end-to-end",
    "Reduces manual effort by 90% compared to baseline",
    "Scales horizontally with worker pool architecture",
]

for i, bullet in enumerate(bullets):
    if i == 0:
        p = left_tf.paragraphs[0]
    else:
        p = left_tf.add_paragraph()
    p.space_before = Pt(12)
    p.space_after = Pt(4)

    run_marker = p.add_run()
    run_marker.text = "\u25a0  "
    run_marker.font.size = Pt(14)
    run_marker.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    run_marker.font.name = BRAND_FONT_BODY

    run_text = p.add_run()
    run_text.text = bullet
    run_text.font.size = Pt(18)
    run_text.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    run_text.font.name = BRAND_FONT_BODY

# ── RIGHT COLUMN: Image placeholder ─────────────────────────────────
# Replace this with slide.shapes.add_picture() for real images
img_placeholder = slide.shapes.add_shape(1, Inches(7.2), Inches(2.0), Inches(5.2), Inches(4.8))
img_placeholder.fill.solid()
img_placeholder.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
img_placeholder.line.color.rgb = hex_to_rgb("#E5E7EB")
img_placeholder.line.width = Pt(1)

# Placeholder label
img_label = slide.shapes.add_textbox(Inches(7.2), Inches(4.0), Inches(5.2), Inches(0.8))
img_tf = img_label.text_frame
img_tf.word_wrap = True
img_p = img_tf.paragraphs[0]
img_p.text = "[ Image / Diagram ]"
img_p.font.size = Pt(16)
img_p.font.color.rgb = hex_to_rgb(BRAND_TEXT_LIGHT)
img_p.font.name = BRAND_FONT_BODY
img_p.alignment = PP_ALIGN.CENTER

# ── Speaker notes ────────────────────────────────────────────────────
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Speaker notes: Walk through the left-column points while referencing the visual on the right."
)

prs.save("two-column-slide.pptx")
print("Created: two-column-slide.pptx")
