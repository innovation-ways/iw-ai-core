# /// script
# name = "quote-slide"
# purpose = "Testimonial or impactful quote with large text and attribution"
# best_for = "Client testimonials, impactful statistics with context, mission statements"
# avoid_when = "Multiple data points (use stats-slide), technical content (use content-slide)"
# instructions = "Replace BRAND_* variables with brand.json values. Quote should be 20-40 words. Attribution line: name + role/company."
# max_quote_words = 40
# requires = ["python-pptx"]
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BRAND_PRIMARY = "#1B2A4A"
BRAND_SECONDARY = "#2E86AB"
BRAND_ACCENT = "#F18F01"
BRAND_BG = "#FFFFFF"
BRAND_BG_ALT = "#F5F7FA"
BRAND_TEXT = "#1A1A2E"
BRAND_TEXT_LIGHT = "#6B7280"
BRAND_TEXT_ON_PRIMARY = "#FFFFFF"
BRAND_FONT_HEADING = "Inter"
BRAND_FONT_BODY = "Inter"


def hex_to_rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Background ───────────────────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)

# ── Left accent bar ─────────────────────────────────────────────────
bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.25), Inches(7.5))
bar.fill.solid()
bar.fill.fore_color.rgb = hex_to_rgb(BRAND_SECONDARY)
bar.line.fill.background()

# ── Large opening quote mark ─────────────────────────────────────────
quote_mark = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(2.0), Inches(1.5))
qm_tf = quote_mark.text_frame
qm_p = qm_tf.paragraphs[0]
qm_p.text = "\u201c"
qm_p.font.size = Pt(120)
qm_p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
qm_p.font.name = BRAND_FONT_HEADING
qm_p.alignment = PP_ALIGN.LEFT

# ── Quote text ───────────────────────────────────────────────────────
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.5), Inches(9.5), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The system reduced our episode production cost from $2.40 to $0.08 while improving audio quality across the board."
p.font.size = Pt(28)
p.font.color.rgb = hex_to_rgb(BRAND_PRIMARY)
p.font.name = BRAND_FONT_HEADING
p.font.italic = True
p.alignment = PP_ALIGN.LEFT
p.line_spacing = Pt(40)

# ── Attribution ──────────────────────────────────────────────────────
attr_box = slide.shapes.add_textbox(Inches(2.0), Inches(5.3), Inches(9.5), Inches(0.8))
attr_tf = attr_box.text_frame
attr_tf.word_wrap = True

# Dash + name
attr_p = attr_tf.paragraphs[0]
run_dash = attr_p.add_run()
run_dash.text = "\u2014  "
run_dash.font.size = Pt(18)
run_dash.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
run_dash.font.name = BRAND_FONT_BODY

run_name = attr_p.add_run()
run_name.text = "Name Surname, Role at Company"
run_name.font.size = Pt(18)
run_name.font.color.rgb = hex_to_rgb(BRAND_TEXT_LIGHT)
run_name.font.name = BRAND_FONT_BODY

# ── Speaker notes ────────────────────────────────────────────────────
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = "Speaker notes: Pause and let the quote resonate. Add context about the relationship with the quoted person."

prs.save("quote-slide.pptx")
print("Created: quote-slide.pptx")
