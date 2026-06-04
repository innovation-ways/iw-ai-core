# /// script
# name = "section-break-slide"
# purpose = "Chapter divider with accent background and large centered title"
# best_for = "Separating major presentation sections, topic transitions, agenda markers"
# avoid_when = "First slide (use title-slide), last slide (use closing-slide)"
# instructions = "Replace BRAND_* variables with brand.json values. Section title should be 2-5 words. Optional subtitle for context."
# max_title_words = 5
# requires = ["python-pptx"]
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BRAND_PRIMARY = "#1B2A4A"
BRAND_SECONDARY = "#2E86AB"
BRAND_ACCENT = "#F18F01"
BRAND_TEXT_ON_PRIMARY = "#FFFFFF"
BRAND_FONT_HEADING = "Inter"
BRAND_FONT_BODY = "Inter"


def hex_to_rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Background — accent color ────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = hex_to_rgb(BRAND_SECONDARY)

# ── Top decorative bar ───────────────────────────────────────────────
top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
top_bar.line.fill.background()

# ── Section number / label ───────────────────────────────────────────
num_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(0.8))
num_tf = num_box.text_frame
num_p = num_tf.paragraphs[0]
num_p.text = "SECTION 01"
num_p.font.size = Pt(16)
num_p.font.bold = True
num_p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
num_p.font.name = BRAND_FONT_BODY
num_p.alignment = PP_ALIGN.LEFT

# ── Section title ────────────────────────────────────────────────────
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.333), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Section Title"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = hex_to_rgb(BRAND_TEXT_ON_PRIMARY)
p.font.name = BRAND_FONT_HEADING
p.alignment = PP_ALIGN.LEFT

# ── Optional subtitle ────────────────────────────────────────────────
sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10.333), Inches(0.8))
sub_tf = sub_box.text_frame
sub_tf.word_wrap = True
sub_p = sub_tf.paragraphs[0]
sub_p.text = "Brief context for this section"
sub_p.font.size = Pt(20)
sub_p.font.color.rgb = RGBColor(255, 255, 255)
sub_p.font.name = BRAND_FONT_BODY
sub_p.alignment = PP_ALIGN.LEFT

# ── Bottom decorative bar ────────────────────────────────────────────
bottom_bar = slide.shapes.add_shape(1, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15))
bottom_bar.fill.solid()
bottom_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_PRIMARY)
bottom_bar.line.fill.background()

# ── Speaker notes ────────────────────────────────────────────────────
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = "Speaker notes: Use this transition to check in with the audience. Briefly preview what this section covers."

prs.save("section-break-slide.pptx")
print("Created: section-break-slide.pptx")
