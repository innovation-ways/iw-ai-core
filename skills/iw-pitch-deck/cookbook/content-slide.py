# /// script
# name = "content-slide"
# purpose = "Bullet point slide with left accent bar, title, and 3-6 square-marked bullets"
# best_for = "Key points, feature lists, process steps, requirements"
# avoid_when = "Data-heavy content (use stats-slide), visual comparisons (use two-column-slide)"
# instructions = "Replace BRAND_* variables with brand.json values. Keep bullets to 3-6 items, each under 15 words. Use square markers for visual consistency."
# max_bullets = 6
# max_bullet_chars = 80
# requires = ["python-pptx"]
# ///

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand placeholders ───────────────────────────────────────────────
BRAND_PRIMARY = "#1B2A4A"
BRAND_SECONDARY = "#2E86AB"
BRAND_ACCENT = "#F18F01"
BRAND_BG = "#FFFFFF"
BRAND_BG_ALT = "#F5F7FA"
BRAND_TEXT = "#1A1A2E"
BRAND_TEXT_LIGHT = "#6B7280"
BRAND_FONT_HEADING = "Inter"
BRAND_FONT_BODY = "Inter"


def hex_to_rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Background ───────────────────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

# ── Left accent bar ─────────────────────────────────────────────────
bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.25), Inches(7.5))
bar.fill.solid()
bar.fill.fore_color.rgb = hex_to_rgb(BRAND_PRIMARY)
bar.line.fill.background()

# ── Slide title ──────────────────────────────────────────────────────
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Slide Title Here"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = hex_to_rgb(BRAND_PRIMARY)
p.font.name = BRAND_FONT_HEADING
p.alignment = PP_ALIGN.LEFT

# ── Title underline ──────────────────────────────────────────────────
line = slide.shapes.add_shape(1, Inches(1.0), Inches(1.55), Inches(2.0), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
line.line.fill.background()

# ── Bullet points ────────────────────────────────────────────────────
bullets = [
    "First key point with a concise message",
    "Second point expanding on the value proposition",
    "Third point with supporting evidence or metric",
    "Fourth point covering implementation detail",
]

txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.0), Inches(4.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True

for i, bullet in enumerate(bullets):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.space_before = Pt(12)
    p.space_after = Pt(4)

    # Square marker
    run_marker = p.add_run()
    run_marker.text = "\u25a0  "
    run_marker.font.size = Pt(14)
    run_marker.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    run_marker.font.name = BRAND_FONT_BODY

    # Bullet text
    run_text = p.add_run()
    run_text.text = bullet
    run_text.font.size = Pt(20)
    run_text.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    run_text.font.name = BRAND_FONT_BODY

# ── Speaker notes ────────────────────────────────────────────────────
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Speaker notes: Expand on each bullet with examples and supporting data."
)

prs.save("content-slide.pptx")
print("Created: content-slide.pptx")
